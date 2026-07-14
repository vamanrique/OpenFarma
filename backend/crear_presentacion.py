"""
crear_presentacion.py — Genera RECURSOS/Presentacion.pptx para el concurso
Datos al Ecosistema 2026: IA para Colombia — Categoría Avanzado

Ejecutar desde la raíz del repo:
  python backend/crear_presentacion.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from pathlib import Path

# ── Paleta de colores OpenFarma ────────────────────────────────────────────────
AZUL_OSCURO   = RGBColor(0x0D, 0x2B, 0x55)   # #0D2B55 — fondo principal
AZUL_MEDIO    = RGBColor(0x1A, 0x56, 0x8C)   # #1A568C — acentos
VERDE_SALUD   = RGBColor(0x00, 0xA8, 0x7E)   # #00A87E — highlight positivo
ROJO_ALERTA   = RGBColor(0xE8, 0x3A, 0x3A)   # #E83A3A — alertas/riesgo
GRIS_CLARO    = RGBColor(0xF2, 0xF6, 0xFC)   # #F2F6FC — fondo claro
BLANCO        = RGBColor(0xFF, 0xFF, 0xFF)
AMARILLO      = RGBColor(0xF5, 0xA6, 0x23)   # #F5A623 — riesgo medio
GRIS_TEXTO    = RGBColor(0x44, 0x4D, 0x5A)   # #444D5A — texto gris

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def prs_nueva():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def fondo_color(slide, color: RGBColor):
    """Rellena el fondo completo de un slide con un color sólido."""
    from pptx.util import Emu
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=Pt(14), bold=False, color=BLANCO,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_slide(slide, title, bullets, title_color=BLANCO,
                     bullet_color=BLANCO, accent_color=VERDE_SALUD,
                     bg_color=AZUL_OSCURO):
    """Slide estándar con título + lista de bullets."""
    fondo_color(slide, bg_color)
    # Barra superior
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=accent_color)
    # Título
    add_text_box(slide, title, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9),
                 font_size=Pt(28), bold=True, color=title_color, align=PP_ALIGN.LEFT)
    # Separador
    add_rect(slide, Inches(0.5), Inches(1.05), Inches(2.5), Inches(0.05), fill_color=accent_color)
    # Bullets
    y = Inches(1.25)
    for bullet in bullets:
        if isinstance(bullet, tuple):
            icon, texto = bullet
        else:
            icon, texto = "▸", bullet
        add_text_box(slide, icon, Inches(0.5), y, Inches(0.4), Inches(0.45),
                     font_size=Pt(16), bold=True, color=accent_color)
        add_text_box(slide, texto, Inches(1.0), y, Inches(11.8), Inches(0.45),
                     font_size=Pt(16), color=bullet_color)
        y += Inches(0.52)
    # Logo esquina inferior derecha
    add_text_box(slide, "OpenFarma  |  Datos al Ecosistema 2026",
                 Inches(8), Inches(7.1), Inches(5), Inches(0.35),
                 font_size=Pt(9), color=RGBColor(0x88, 0x99, 0xAA),
                 align=PP_ALIGN.RIGHT)


def add_two_col_slide(slide, title, col1_title, col1_items,
                      col2_title, col2_items, accent=VERDE_SALUD):
    fondo_color(slide, AZUL_OSCURO)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=accent)
    add_text_box(slide, title, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9),
                 font_size=Pt(28), bold=True, color=BLANCO)
    add_rect(slide, Inches(0.5), Inches(1.05), Inches(2.5), Inches(0.05), fill_color=accent)

    # Columna 1
    add_rect(slide, Inches(0.4), Inches(1.2), Inches(5.9), Inches(5.8),
             fill_color=RGBColor(0x10, 0x35, 0x65))
    add_text_box(slide, col1_title, Inches(0.6), Inches(1.35), Inches(5.5), Inches(0.5),
                 font_size=Pt(14), bold=True, color=accent)
    y = Inches(1.9)
    for item in col1_items:
        add_text_box(slide, f"▸  {item}", Inches(0.7), y, Inches(5.3), Inches(0.5),
                     font_size=Pt(13), color=BLANCO)
        y += Inches(0.52)

    # Columna 2
    add_rect(slide, Inches(6.9), Inches(1.2), Inches(5.9), Inches(5.8),
             fill_color=RGBColor(0x10, 0x35, 0x65))
    add_text_box(slide, col2_title, Inches(7.1), Inches(1.35), Inches(5.5), Inches(0.5),
                 font_size=Pt(14), bold=True, color=accent)
    y = Inches(1.9)
    for item in col2_items:
        add_text_box(slide, f"▸  {item}", Inches(7.2), y, Inches(5.3), Inches(0.5),
                     font_size=Pt(13), color=BLANCO)
        y += Inches(0.52)

    add_text_box(slide, "OpenFarma  |  Datos al Ecosistema 2026",
                 Inches(8), Inches(7.1), Inches(5), Inches(0.35),
                 font_size=Pt(9), color=RGBColor(0x88, 0x99, 0xAA),
                 align=PP_ALIGN.RIGHT)


def add_metric_box(slide, left, top, width, height,
                   label, value, sublabel="", accent=VERDE_SALUD):
    add_rect(slide, left, top, width, height,
             fill_color=RGBColor(0x10, 0x35, 0x65))
    add_rect(slide, left, top, width, Inches(0.07), fill_color=accent)
    add_text_box(slide, value, left, top + Inches(0.15), width, Inches(0.8),
                 font_size=Pt(32), bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, left, top + Inches(0.9), width, Inches(0.5),
                 font_size=Pt(12), bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    if sublabel:
        add_text_box(slide, sublabel, left, top + Inches(1.35), width, Inches(0.35),
                     font_size=Pt(10), color=RGBColor(0x88, 0x99, 0xAA),
                     align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_portada(prs):
    """Slide 1: Portada."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fondo_color(slide, AZUL_OSCURO)

    # Banda lateral izquierda
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, fill_color=VERDE_SALUD)

    # Banda inferior
    add_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), fill_color=RGBColor(0x08, 0x1A, 0x35))

    # Nombre del proyecto — grande
    add_text_box(slide, "OpenFarma", Inches(0.5), Inches(0.6), Inches(9), Inches(1.6),
                 font_size=Pt(72), bold=True, color=BLANCO)

    # Línea verde decorativa
    add_rect(slide, Inches(0.5), Inches(2.1), Inches(3), Inches(0.07), fill_color=VERDE_SALUD)

    # Tagline
    add_text_box(slide,
                 "Sistema de Alerta Temprana de\nDesabastecimiento Farmacéutico",
                 Inches(0.5), Inches(2.25), Inches(9), Inches(1.2),
                 font_size=Pt(24), bold=False, color=GRIS_CLARO)

    # Descripción
    add_text_box(slide,
                 "Predicción con IA · Datos CUM + INVIMA · Canal ciudadano",
                 Inches(0.5), Inches(3.55), Inches(9), Inches(0.6),
                 font_size=Pt(16), color=RGBColor(0x88, 0xCC, 0xEE))

    # Badge concurso (caja verde)
    add_rect(slide, Inches(0.5), Inches(4.3), Inches(4.5), Inches(0.65),
             fill_color=VERDE_SALUD)
    add_text_box(slide, "Datos al Ecosistema 2026: IA para Colombia",
                 Inches(0.55), Inches(4.35), Inches(4.4), Inches(0.55),
                 font_size=Pt(13), bold=True, color=BLANCO)

    # Badge categoría
    add_rect(slide, Inches(5.1), Inches(4.3), Inches(2.2), Inches(0.65),
             fill_color=AZUL_MEDIO)
    add_text_box(slide, "Categoría: Avanzado",
                 Inches(5.15), Inches(4.35), Inches(2.1), Inches(0.55),
                 font_size=Pt(13), bold=True, color=BLANCO)

    # Equipo — BORRADOR
    add_rect(slide, Inches(0.5), Inches(5.15), Inches(6.8), Inches(1.3),
             fill_color=RGBColor(0x10, 0x35, 0x65))
    add_rect(slide, Inches(0.5), Inches(5.15), Inches(6.8), Inches(0.07),
             fill_color=AMARILLO)
    add_text_box(slide, "⚠  EQUIPO — BORRADOR (pendiente inscripción)",
                 Inches(0.6), Inches(5.22), Inches(6.5), Inches(0.4),
                 font_size=Pt(11), bold=True, color=AMARILLO)
    add_text_box(slide,
                 "Líder: [NOMBRE]  ·  Analista de Datos: [NOMBRE]  ·  Desarrollador: [NOMBRE]",
                 Inches(0.6), Inches(5.62), Inches(6.5), Inches(0.35),
                 font_size=Pt(11), color=RGBColor(0xBB, 0xCC, 0xDD))
    add_text_box(slide, "Contacto: vamanrique@gmail.com",
                 Inches(0.6), Inches(5.95), Inches(6.5), Inches(0.35),
                 font_size=Pt(11), color=RGBColor(0x88, 0xCC, 0xEE))

    # Métricas a la derecha (mini dashboard)
    add_rect(slide, Inches(10.0), Inches(0.5), Inches(3.0), Inches(6.0),
             fill_color=RGBColor(0x08, 0x1A, 0x35))
    for i, (val, lbl, sub) in enumerate([
        ("52,830", "Medicamentos CUM", "normalizados"),
        ("9,795",  "Alertas INVIMA",   "17 meses"),
        ("0.8374", "ROC-AUC",          "split temporal"),
        ("3,204",  "Grupos terapéuticos", "equivalencia"),
    ]):
        y = Inches(0.6 + i * 1.35)
        add_rect(slide, Inches(10.1), y, Inches(2.8), Inches(1.2),
                 fill_color=RGBColor(0x10, 0x35, 0x65))
        add_text_box(slide, val, Inches(10.1), y + Inches(0.05), Inches(2.8), Inches(0.6),
                     font_size=Pt(26), bold=True, color=VERDE_SALUD, align=PP_ALIGN.CENTER)
        add_text_box(slide, lbl, Inches(10.1), y + Inches(0.6), Inches(2.8), Inches(0.35),
                     font_size=Pt(10), bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        add_text_box(slide, sub, Inches(10.1), y + Inches(0.9), Inches(2.8), Inches(0.25),
                     font_size=Pt(9), color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)

    # URL demo
    add_text_box(slide, "🌐  openfarma-production.up.railway.app",
                 Inches(0.5), Inches(6.85), Inches(7), Inches(0.4),
                 font_size=Pt(11), color=RGBColor(0x88, 0xCC, 0xEE))
    add_text_box(slide, "github.com/vamanrique/OpenFarma",
                 Inches(7.5), Inches(6.85), Inches(5.5), Inches(0.4),
                 font_size=Pt(11), color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.RIGHT)


def slide_02_problema(prs):
    """Slide 2: El Problema."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_color(slide, AZUL_OSCURO)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=ROJO_ALERTA)

    add_text_box(slide, "El Problema: Desabastecimiento Farmacéutico",
                 Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9),
                 font_size=Pt(28), bold=True, color=BLANCO)
    add_rect(slide, Inches(0.5), Inches(1.05), Inches(3), Inches(0.05), fill_color=ROJO_ALERTA)

    add_text_box(slide,
                 "Cuando INVIMA publica una alerta de desabastecimiento,\nel medicamento ya lleva semanas sin llegar a las farmacias.",
                 Inches(0.5), Inches(1.2), Inches(8.5), Inches(0.9),
                 font_size=Pt(17), italic=True, color=RGBColor(0xDD, 0xEE, 0xFF))

    # 3 cajas de cifras impactantes
    for i, (num, lbl, desc) in enumerate([
        ("9,795", "alertas INVIMA", "en 17 meses (ene 2025 – may 2026)"),
        ("5+",    "semanas de retraso", "entre escasez real y alerta oficial"),
        ("0",     "sistemas preventivos", "públicos integrados en Colombia"),
    ]):
        x = Inches(0.4 + i * 4.3)
        add_rect(slide, x, Inches(2.2), Inches(4.0), Inches(2.0),
                 fill_color=RGBColor(0x10, 0x35, 0x65))
        add_rect(slide, x, Inches(2.2), Inches(4.0), Inches(0.08), fill_color=ROJO_ALERTA)
        add_text_box(slide, num, x, Inches(2.35), Inches(4.0), Inches(0.85),
                     font_size=Pt(40), bold=True, color=ROJO_ALERTA, align=PP_ALIGN.CENTER)
        add_text_box(slide, lbl, x, Inches(3.1), Inches(4.0), Inches(0.45),
                     font_size=Pt(13), bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x, Inches(3.5), Inches(4.0), Inches(0.55),
                     font_size=Pt(10), color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)

    add_text_box(slide, "El impacto:", Inches(0.5), Inches(4.4), Inches(2), Inches(0.45),
                 font_size=Pt(14), bold=True, color=AMARILLO)

    for i, txt in enumerate([
        "Pacientes con enfermedades crónicas interrumpen tratamientos",
        "IPS y clínicas no pueden anticipar compras de emergencia",
        "INVIMA reacciona cuando la escasez ya es un hecho consumado",
        "No existe señal ciudadana integrada con los datos regulatorios",
    ]):
        add_text_box(slide, f"▸  {txt}", Inches(0.5), Inches(4.85 + i * 0.48),
                     Inches(12.4), Inches(0.45), font_size=Pt(14), color=BLANCO)

    add_text_box(slide, "OpenFarma  |  Datos al Ecosistema 2026",
                 Inches(8), Inches(7.1), Inches(5), Inches(0.35),
                 font_size=Pt(9), color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.RIGHT)


def slide_03_datos(prs):
    """Slide 3: Los Datos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_color(slide, AZUL_OSCURO)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=VERDE_SALUD)

    add_text_box(slide, "Los Datos: Fuentes Públicas Integradas por Primera Vez",
                 Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9),
                 font_size=Pt(26), bold=True, color=BLANCO)
    add_rect(slide, Inches(0.5), Inches(1.05), Inches(3), Inches(0.05), fill_color=VERDE_SALUD)

    # 4 fuentes de datos
    sources = [
        ("CUM Activos",        "52,830 presentaciones",  "datos.gov.co · i7cb-raxc",    VERDE_SALUD),
        ("CUM Renovación",     "~8,000 registros",        "datos.gov.co · vgr4-gemg",    AZUL_MEDIO),
        ("Alertas INVIMA",     "9,795 entradas limpias",  "PDFs portal INVIMA · 17 meses", AMARILLO),
        ("Reportes ciudadanos","Canal propio",             "Formulario OpenFarma · activo",ROJO_ALERTA),
    ]
    for i, (title, metric, source, color) in enumerate(sources):
        col = i % 2
        row = i // 2
        x = Inches(0.4 + col * 6.5)
        y = Inches(1.25 + row * 2.7)
        add_rect(slide, x, y, Inches(6.1), Inches(2.4), fill_color=RGBColor(0x10, 0x35, 0x65))
        add_rect(slide, x, y, Inches(6.1), Inches(0.08), fill_color=color)
        add_text_box(slide, title, x + Inches(0.15), y + Inches(0.15),
                     Inches(5.8), Inches(0.5), font_size=Pt(16), bold=True, color=color)
        add_text_box(slide, metric, x + Inches(0.15), y + Inches(0.65),
                     Inches(5.8), Inches(0.7), font_size=Pt(28), bold=True, color=BLANCO)
        add_text_box(slide, source, x + Inches(0.15), y + Inches(1.4),
                     Inches(5.8), Inches(0.4), font_size=Pt(11),
                     color=RGBColor(0x88, 0x99, 0xAA))

    # Nota diferenciadora
    add_rect(slide, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.55),
             fill_color=RGBColor(0x00, 0x40, 0x30))
    add_text_box(slide,
                 "✓  Primera integración pública completa CUM + INVIMA + Reportes Ciudadanos en Colombia",
                 Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.45),
                 font_size=Pt(13), bold=True, color=VERDE_SALUD)


def slide_04_arquitectura(prs):
    """Slide 4: Arquitectura de la Solución."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_color(slide, AZUL_OSCURO)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=VERDE_SALUD)

    add_text_box(slide, "Arquitectura: Tres Componentes Integrados",
                 Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9),
                 font_size=Pt(28), bold=True, color=BLANCO)
    add_rect(slide, Inches(0.5), Inches(1.05), Inches(3), Inches(0.05), fill_color=VERDE_SALUD)

    # Capa Ciudadano
    add_rect(slide, Inches(0.3), Inches(1.2), Inches(2.2), Inches(5.8),
             fill_color=RGBColor(0x08, 0x22, 0x45))
    add_rect(slide, Inches(0.3), Inches(1.2), Inches(2.2), Inches(0.08), fill_color=VERDE_SALUD)
    add_text_box(slide, "CIUDADANO", Inches(0.3), Inches(1.3), Inches(2.2), Inches(0.5),
                 font_size=Pt(11), bold=True, color=VERDE_SALUD, align=PP_ALIGN.CENTER)
    for item in ["Busca medicamento", "Ve alternativas", "Consulta riesgo ML", "Reporta escasez", "Ve mapa Colombia"]:
        y_off = Inches(1.9 + ["Busca medicamento","Ve alternativas","Consulta riesgo ML","Reporta escasez","Ve mapa Colombia"].index(item) * 0.9)
        add_text_box(slide, item, Inches(0.35), y_off, Inches(2.1), Inches(0.8),
                     font_size=Pt(11), color=BLANCO, align=PP_ALIGN.CENTER)

    # Flecha
    add_text_box(slide, "→", Inches(2.55), Inches(3.8), Inches(0.5), Inches(0.5),
                 font_size=Pt(24), bold=True, color=VERDE_SALUD)

    # Frontend React
    add_rect(slide, Inches(3.1), Inches(1.2), Inches(2.2), Inches(5.8),
             fill_color=RGBColor(0x08, 0x22, 0x45))
    add_rect(slide, Inches(3.1), Inches(1.2), Inches(2.2), Inches(0.08), fill_color=AZUL_MEDIO)
    add_text_box(slide, "REACT FRONTEND", Inches(3.1), Inches(1.3), Inches(2.2), Inches(0.5),
                 font_size=Pt(11), bold=True, color=AZUL_MEDIO, align=PP_ALIGN.CENTER)
    for item in ["React 19 + Vite", "Tailwind CSS", "Leaflet Maps", "Recharts", "ARIA a11y"]:
        y_off = Inches(1.9 + ["React 19 + Vite","Tailwind CSS","Leaflet Maps","Recharts","ARIA a11y"].index(item) * 0.9)
        add_text_box(slide, item, Inches(3.15), y_off, Inches(2.1), Inches(0.8),
                     font_size=Pt(11), color=BLANCO, align=PP_ALIGN.CENTER)

    add_text_box(slide, "→", Inches(5.35), Inches(3.8), Inches(0.5), Inches(0.5),
                 font_size=Pt(24), bold=True, color=AZUL_MEDIO)

    # FastAPI Backend
    add_rect(slide, Inches(5.9), Inches(1.2), Inches(2.2), Inches(5.8),
             fill_color=RGBColor(0x08, 0x22, 0x45))
    add_rect(slide, Inches(5.9), Inches(1.2), Inches(2.2), Inches(0.08), fill_color=AMARILLO)
    add_text_box(slide, "FASTAPI BACKEND", Inches(5.9), Inches(1.3), Inches(2.2), Inches(0.5),
                 font_size=Pt(11), bold=True, color=AMARILLO, align=PP_ALIGN.CENTER)
    for item in ["Búsqueda CUM", "Alternativas", "Predicciones ML", "Reportes API", "INVIMA Cache"]:
        y_off = Inches(1.9 + ["Búsqueda CUM","Alternativas","Predicciones ML","Reportes API","INVIMA Cache"].index(item) * 0.9)
        add_text_box(slide, item, Inches(5.95), y_off, Inches(2.1), Inches(0.8),
                     font_size=Pt(11), color=BLANCO, align=PP_ALIGN.CENTER)

    add_text_box(slide, "→", Inches(8.15), Inches(3.8), Inches(0.5), Inches(0.5),
                 font_size=Pt(24), bold=True, color=AMARILLO)

    # Datos + ML
    add_rect(slide, Inches(8.7), Inches(1.2), Inches(2.2), Inches(5.8),
             fill_color=RGBColor(0x08, 0x22, 0x45))
    add_rect(slide, Inches(8.7), Inches(1.2), Inches(2.2), Inches(0.08), fill_color=ROJO_ALERTA)
    add_text_box(slide, "DATOS + ML", Inches(8.7), Inches(1.3), Inches(2.2), Inches(0.5),
                 font_size=Pt(11), bold=True, color=ROJO_ALERTA, align=PP_ALIGN.CENTER)
    for item in ["SQLite 52K CUM", "3,204 grupos", "INVIMA 17m", "RF + Calibrado", "Bias Tests"]:
        y_off = Inches(1.9 + ["SQLite 52K CUM","3,204 grupos","INVIMA 17m","RF + Calibrado","Bias Tests"].index(item) * 0.9)
        add_text_box(slide, item, Inches(8.75), y_off, Inches(2.1), Inches(0.8),
                     font_size=Pt(11), color=BLANCO, align=PP_ALIGN.CENTER)

    add_text_box(slide, "→", Inches(10.95), Inches(3.8), Inches(0.5), Inches(0.5),
                 font_size=Pt(24), bold=True, color=ROJO_ALERTA)

    # Alerta
    add_rect(slide, Inches(11.5), Inches(1.2), Inches(1.5), Inches(5.8),
             fill_color=RGBColor(0x3A, 0x10, 0x10))
    add_rect(slide, Inches(11.5), Inches(1.2), Inches(1.5), Inches(0.08), fill_color=ROJO_ALERTA)
    add_text_box(slide, "ALERTA\nTEMPRANA", Inches(11.5), Inches(3.5), Inches(1.5), Inches(1.2),
                 font_size=Pt(13), bold=True, color=ROJO_ALERTA, align=PP_ALIGN.CENTER)

    # Railway badge
    add_text_box(slide, "Deploy: Railway · Auto-deploy desde GitHub main",
                 Inches(0.5), Inches(7.1), Inches(7), Inches(0.35),
                 font_size=Pt(9), color=RGBColor(0x88, 0x99, 0xAA))
    add_text_box(slide, "OpenFarma  |  Datos al Ecosistema 2026",
                 Inches(8), Inches(7.1), Inches(5), Inches(0.35),
                 font_size=Pt(9), color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.RIGHT)


def slide_05_pipeline(prs):
    """Slide 5: Pipeline ETL & Normalización."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(
        slide,
        "Pipeline ETL: Calidad Farmacéutica de Alta Precisión",
        bullets=[
            ("▸", "52,830 medicamentos del CUM descargados vía Socrata API (datos.gov.co)"),
            ("▸", "Normalización INN: 420 reglas de sinonimia + 50 patrones de sal farmacéutica"),
            ("▸", "3,204 grupos de equivalencia terapéutica construidos con DeepSeek + reglas"),
            ("▸", "105 rondas de auditoría INN: radiofármacos, vacunas, biológicos, hemoderivados"),
            ("▸", "9,795 alertas INVIMA parseadas de PDFs mensuales (17 meses, ene 2025–may 2026)"),
            ("▸", "0 duplicados · 0 NULL concentración · 0 mismatches DCI tras auditoría completa"),
        ],
        accent_color=VERDE_SALUD,
    )
    # Métricas de calidad en fila
    y = Inches(5.4)
    for i, (val, lbl) in enumerate([
        ("100%", "DCIs normalizados"),
        ("105",  "Rondas auditoría"),
        ("0",    "Duplicados"),
        ("17",   "Meses INVIMA"),
    ]):
        x = Inches(0.4 + i * 3.2)
        add_rect(slide, x, y, Inches(3.0), Inches(1.5),
                 fill_color=RGBColor(0x10, 0x35, 0x65))
        add_text_box(slide, val, x, y + Inches(0.1), Inches(3.0), Inches(0.8),
                     font_size=Pt(30), bold=True, color=VERDE_SALUD, align=PP_ALIGN.CENTER)
        add_text_box(slide, lbl, x, y + Inches(0.85), Inches(3.0), Inches(0.5),
                     font_size=Pt(11), color=BLANCO, align=PP_ALIGN.CENTER)
    add_text_box(slide,
                 "Fallback local: si datos.gov.co no responde, los 52,830 productos siguen disponibles",
                 Inches(0.5), Inches(7.0), Inches(12), Inches(0.35),
                 font_size=Pt(10), italic=True, color=RGBColor(0x88, 0x99, 0xAA))


def slide_06_modelo(prs):
    """Slide 6: Modelo ML."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_color(slide, AZUL_OSCURO)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=ROJO_ALERTA)

    add_text_box(slide, "Modelo Predictivo: Integridad Estadística Rigurosa",
                 Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9),
                 font_size=Pt(26), bold=True, color=BLANCO)
    add_rect(slide, Inches(0.5), Inches(1.05), Inches(3), Inches(0.05), fill_color=ROJO_ALERTA)

    # Métricas principales
    add_metric_box(slide, Inches(0.4), Inches(1.2), Inches(3.1), Inches(2.0),
                   "ROC-AUC", "0.8374", "split temporal honesto", VERDE_SALUD)
    add_metric_box(slide, Inches(3.65), Inches(1.2), Inches(3.1), Inches(2.0),
                   "Avg Precision", "0.1707", "1.6% positivos reales", AMARILLO)
    add_metric_box(slide, Inches(6.9), Inches(1.2), Inches(3.1), Inches(2.0),
                   "Meses test", "3", "mar–may 2026 (nunca vistos)", AZUL_MEDIO)
    add_metric_box(slide, Inches(10.15), Inches(1.2), Inches(2.8), Inches(2.0),
                   "Features", "15", "10 CUM + 5 INVIMA", ROJO_ALERTA)

    # Split temporal explicado
    add_rect(slide, Inches(0.4), Inches(3.4), Inches(6.0), Inches(3.8),
             fill_color=RGBColor(0x10, 0x35, 0x65))
    add_text_box(slide, "Split Temporal (sin Data Leakage)",
                 Inches(0.5), Inches(3.5), Inches(5.8), Inches(0.5),
                 font_size=Pt(14), bold=True, color=VERDE_SALUD)
    add_text_box(slide,
                 "▸  Train: ene 2025 – feb 2026 (14 meses)\n"
                 "▸  Test:  mar – may 2026 (3 meses, nunca vistos)\n"
                 "▸  Con split aleatorio: ROC-AUC era 1.000 (data leakage)\n"
                 "▸  Con split temporal: ROC-AUC 0.8374 (honesto)\n"
                 "▸  Modelo producción: reentrenado en todos los datos",
                 Inches(0.5), Inches(4.0), Inches(5.8), Inches(3.0),
                 font_size=Pt(13), color=BLANCO)

    # Top features
    add_rect(slide, Inches(6.6), Inches(3.4), Inches(6.5), Inches(3.8),
             fill_color=RGBColor(0x10, 0x35, 0x65))
    add_text_box(slide, "Top Features por Importancia",
                 Inches(6.7), Inches(3.5), Inches(6.3), Inches(0.5),
                 font_size=Pt(14), bold=True, color=ROJO_ALERTA)
    features = [
        ("invima_sev_actual",         "27.5%", ROJO_ALERTA),
        ("invima_peor_sev_hist",       "21.1%", ROJO_ALERTA),
        ("invima_meses_monitoreado",   "12.9%", AMARILLO),
        ("tasa_inactivacion_atc5",     "11.6%", AMARILLO),
        ("invima_sev_t3_avg",          "11.5%", AMARILLO),
    ]
    for i, (feat, pct, color) in enumerate(features):
        y = Inches(4.05 + i * 0.6)
        bar_w = float(pct.strip("%")) / 30.0
        add_rect(slide, Inches(6.7), y + Inches(0.05), Inches(bar_w), Inches(0.38),
                 fill_color=RGBColor(0x1A, 0x56, 0x8C))
        add_text_box(slide, feat, Inches(6.7), y, Inches(4.5), Inches(0.45),
                     font_size=Pt(11), color=BLANCO)
        add_text_box(slide, pct, Inches(11.8), y, Inches(1.1), Inches(0.45),
                     font_size=Pt(12), bold=True, color=color, align=PP_ALIGN.RIGHT)

    add_text_box(slide,
                 "Tipo: CalibratedClassifierCV (Platt scaling) + RandomForestClassifier · scikit-learn 1.9.0",
                 Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
                 font_size=Pt(9), color=RGBColor(0x88, 0x99, 0xAA))
    add_text_box(slide, "OpenFarma  |  Datos al Ecosistema 2026",
                 Inches(8), Inches(7.1), Inches(5), Inches(0.35),
                 font_size=Pt(9), color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.RIGHT)


def slide_07_aplicacion(prs):
    """Slide 7: La Aplicación."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(
        slide,
        "La Aplicación: Accesible para Pacientes, Clínicos e INVIMA",
        bullets=[
            ("🔍", "Búsqueda en tiempo real en 52,830 medicamentos CUM (Socrata API + fallback local)"),
            ("💊", "Alternativas terapéuticas en 8 niveles: sustituto directo → clase ATC completa"),
            ("🤖", "Badge de riesgo ML por medicamento: probabilidad 0-100% con nivel Bajo/Medio/Alto/Crítico"),
            ("🗺️",  "Mapa interactivo de Colombia: riesgo por departamento filtrable por nivel"),
            ("📊", "Dashboard de vigilancia: top 20 reportados + spike detector vs. alertas INVIMA"),
            ("🚨", "Señales anticipadas: medicamentos con spike ciudadano sin alerta INVIMA aún"),
        ],
        accent_color=AZUL_MEDIO,
    )
    # URL + tech
    add_rect(slide, Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.75),
             fill_color=RGBColor(0x10, 0x35, 0x65))
    add_text_box(slide,
                 "🌐  https://openfarma-production.up.railway.app  "
                 "·  CI/CD GitHub Actions  ·  16/16 tests verdes  "
                 "·  ARIA accesible  ·  MIT License",
                 Inches(0.55), Inches(6.12), Inches(12.2), Inches(0.55),
                 font_size=Pt(12), color=VERDE_SALUD)
    add_text_box(slide,
                 "Stack: FastAPI · SQLAlchemy · SQLite · React 19 · Vite · Tailwind · Leaflet · Recharts · scikit-learn",
                 Inches(0.5), Inches(7.1), Inches(12), Inches(0.35),
                 font_size=Pt(9), color=RGBColor(0x88, 0x99, 0xAA))


def slide_08_impacto(prs):
    """Slide 8: Impacto y Escalabilidad."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_color(slide, AZUL_OSCURO)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=VERDE_SALUD)

    add_text_box(slide, "Impacto y Escalabilidad",
                 Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9),
                 font_size=Pt(28), bold=True, color=BLANCO)
    add_rect(slide, Inches(0.5), Inches(1.05), Inches(2), Inches(0.05), fill_color=VERDE_SALUD)

    # Columna izquierda: impacto
    add_rect(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(5.4),
             fill_color=RGBColor(0x10, 0x35, 0x65))
    add_text_box(slide, "Impacto Inmediato", Inches(0.55), Inches(1.3),
                 Inches(5.7), Inches(0.5), font_size=Pt(16), bold=True, color=VERDE_SALUD)
    impacts = [
        "52,830 medicamentos monitoreados en tiempo real",
        "Canal ciudadano → señal regulatoria directa a INVIMA",
        "Anticipa desabastecimiento 30 días antes que alertas oficiales",
        "Alternativas terapéuticas: reduce impacto de escasez en pacientes",
        "Código abierto MIT: cualquier entidad puede adoptar el sistema",
    ]
    for i, txt in enumerate(impacts):
        add_text_box(slide, f"✓  {txt}", Inches(0.6), Inches(1.9 + i * 0.75),
                     Inches(5.6), Inches(0.65), font_size=Pt(13), color=BLANCO)

    # Columna derecha: escalabilidad
    add_rect(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.4),
             fill_color=RGBColor(0x10, 0x35, 0x65))
    add_text_box(slide, "Escalabilidad y Roadmap", Inches(6.95), Inches(1.3),
                 Inches(5.7), Inches(0.5), font_size=Pt(16), bold=True, color=AMARILLO)
    roadmap = [
        "🏥  Integración directa API INVIMA (alertas automáticas)",
        "📱  Notificaciones push a IPS y farmacias en riesgo",
        "🌎  Replicable a Ecuador, Perú, Chile (datos similares DIGEMID/ISP)",
        "📈  Re-entreno mensual automático con datos INVIMA nuevos",
        "🤝  Partnership con MinSalud para adopción institucional",
    ]
    for i, txt in enumerate(roadmap):
        add_text_box(slide, txt, Inches(7.0), Inches(1.9 + i * 0.75),
                     Inches(5.5), Inches(0.65), font_size=Pt(13), color=BLANCO)

    # CTA final
    add_rect(slide, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.6),
             fill_color=VERDE_SALUD)
    add_text_box(slide,
                 "⭐  github.com/vamanrique/OpenFarma  ·  "
                 "🌐  openfarma-production.up.railway.app  ·  "
                 "📧  vamanrique@gmail.com",
                 Inches(0.55), Inches(6.82), Inches(12.2), Inches(0.45),
                 font_size=Pt(13), bold=True, color=BLANCO, align=PP_ALIGN.CENTER)


def slide_09_equipo(prs):
    """Slide 9: Equipo — BORRADOR."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo_color(slide, AZUL_OSCURO)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill_color=AMARILLO)

    add_text_box(slide, "El Equipo",
                 Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9),
                 font_size=Pt(28), bold=True, color=BLANCO)
    add_rect(slide, Inches(0.5), Inches(1.05), Inches(2), Inches(0.05), fill_color=AMARILLO)

    # Banner BORRADOR
    add_rect(slide, Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.6),
             fill_color=RGBColor(0x5A, 0x40, 0x00))
    add_text_box(slide,
                 "⚠  BORRADOR — Información de participantes pendiente de confirmación para inscripción oficial",
                 Inches(0.55), Inches(1.2), Inches(12.2), Inches(0.5),
                 font_size=Pt(13), bold=True, color=AMARILLO)

    # 3 placeholders de miembros
    roles = [
        ("Líder de Equipo", "[NOMBRE COMPLETO]", "[Cargo/Institución]", "[Ciudad, Colombia]"),
        ("Analista de Datos / ML", "[NOMBRE COMPLETO]", "[Cargo/Institución]", "[Ciudad, Colombia]"),
        ("Desarrollador Full-Stack", "[NOMBRE COMPLETO]", "[Cargo/Institución]", "[Ciudad, Colombia]"),
    ]
    for i, (role, name, org, city) in enumerate(roles):
        x = Inches(0.4 + i * 4.3)
        add_rect(slide, x, Inches(1.95), Inches(4.0), Inches(4.5),
                 fill_color=RGBColor(0x10, 0x35, 0x65))
        add_rect(slide, x, Inches(1.95), Inches(4.0), Inches(0.08), fill_color=AMARILLO)
        # Avatar placeholder
        add_rect(slide, x + Inches(1.25), Inches(2.1), Inches(1.5), Inches(1.5),
                 fill_color=RGBColor(0x1A, 0x56, 0x8C))
        add_text_box(slide, "👤", x + Inches(1.25), Inches(2.15), Inches(1.5), Inches(1.4),
                     font_size=Pt(40), align=PP_ALIGN.CENTER, color=BLANCO)
        add_text_box(slide, role, x, Inches(3.7), Inches(4.0), Inches(0.5),
                     font_size=Pt(11), bold=True, color=AMARILLO, align=PP_ALIGN.CENTER)
        add_text_box(slide, name, x, Inches(4.15), Inches(4.0), Inches(0.5),
                     font_size=Pt(13), bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        add_text_box(slide, org, x, Inches(4.6), Inches(4.0), Inches(0.4),
                     font_size=Pt(10), color=RGBColor(0xBB, 0xCC, 0xDD),
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, city, x, Inches(4.95), Inches(4.0), Inches(0.4),
                     font_size=Pt(10), color=RGBColor(0x88, 0x99, 0xAA),
                     align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "Requisito concurso: equipos de 2-4 personas · mínimo una mujer · al menos un perfil técnico",
                 Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
                 font_size=Pt(10), italic=True, color=RGBColor(0x88, 0x99, 0xAA))
    add_text_box(slide, "vamanrique@gmail.com",
                 Inches(0.5), Inches(7.05), Inches(5), Inches(0.4),
                 font_size=Pt(11), color=RGBColor(0x88, 0xCC, 0xEE))
    add_text_box(slide, "OpenFarma  |  Datos al Ecosistema 2026",
                 Inches(8), Inches(7.1), Inches(5), Inches(0.35),
                 font_size=Pt(9), color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = prs_nueva()

    print("Generando slides...")
    slide_01_portada(prs)     ; print("  ✓ Slide 1: Portada")
    slide_02_problema(prs)    ; print("  ✓ Slide 2: El Problema")
    slide_03_datos(prs)       ; print("  ✓ Slide 3: Los Datos")
    slide_04_arquitectura(prs); print("  ✓ Slide 4: Arquitectura")
    slide_05_pipeline(prs)    ; print("  ✓ Slide 5: Pipeline ETL")
    slide_06_modelo(prs)      ; print("  ✓ Slide 6: Modelo ML")
    slide_07_aplicacion(prs)  ; print("  ✓ Slide 7: La Aplicación")
    slide_08_impacto(prs)     ; print("  ✓ Slide 8: Impacto y Escalabilidad")
    slide_09_equipo(prs)      ; print("  ✓ Slide 9: Equipo (BORRADOR)")

    out_dir = Path(__file__).parent.parent / "RECURSOS"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "Presentacion.pptx"
    prs.save(str(out_path))
    print(f"\n✅  Guardado: {out_path}")
    print(f"   {len(prs.slides)} slides · {out_path.stat().st_size // 1024} KB")


if __name__ == "__main__":
    main()
